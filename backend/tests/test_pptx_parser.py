"""PPT parser unit tests (W3 D1 F5 — covers parsers/pptx_parser.py).

Uses python-pptx itself to build synthetic .pptx fixtures in tmp dirs since
real Drive Project .pptx samples are blocked on Q2 (PPT-share request from
Chris). Synthetic fixtures cover: title-only slide, body text, table,
picture, speaker notes, multi-slide ordering, malformed file fallback.
"""

from __future__ import annotations

import io
from pathlib import Path

import pytest
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from ingestion.parsers.pptx_parser import PptxParser


def _png_bytes(color: tuple[int, int, int] = (255, 0, 0)) -> bytes:
    """Tiny in-memory PNG for picture-shape fixtures."""
    img = Image.new("RGB", (32, 32), color=color)
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


def _build_pptx(tmp_path: Path, name: str = "fixture.pptx") -> Path:
    """Builds a 2-slide deck with title + body + table + picture + notes."""
    prs = Presentation()
    blank = prs.slide_layouts[5]  # title-only layout

    # Slide 1: title + body text + picture
    s1 = prs.slides.add_slide(blank)
    s1.shapes.title.text = "Sales Overview"
    body = s1.shapes.add_textbox(Inches(1), Inches(2), Inches(6), Inches(2))
    body.text_frame.text = "Q1 revenue grew 12% YoY"
    pic_path = tmp_path / "tmp_pic.png"
    pic_path.write_bytes(_png_bytes())
    s1.shapes.add_picture(str(pic_path), Inches(5), Inches(4), width=Inches(2))
    notes = s1.notes_slide.notes_text_frame
    notes.text = "Speaker note: emphasize APAC growth"

    # Slide 2: title + table
    s2 = prs.slides.add_slide(blank)
    s2.shapes.title.text = "Region Breakdown"
    table_shape = s2.shapes.add_table(
        rows=2, cols=2, left=Inches(1), top=Inches(2), width=Inches(6), height=Inches(2)
    )
    tbl = table_shape.table
    tbl.cell(0, 0).text = "Region"
    tbl.cell(0, 1).text = "Revenue"
    tbl.cell(1, 0).text = "APAC"
    tbl.cell(1, 1).text = "$1.2M"

    out = tmp_path / name
    prs.save(out)
    return out


def test_parse_two_slide_deck_emits_expected_paragraphs(tmp_path: Path) -> None:
    src = _build_pptx(tmp_path)
    result = PptxParser().parse(src)

    assert not result.parse_failed
    assert result.doc_format == "pptx"
    headings = [p for p in result.paragraphs if p.kind == "heading"]
    assert any(p.text == "Slide 1" and p.heading_level == 1 for p in headings)
    assert any(p.text == "Slide 2" and p.heading_level == 1 for p in headings)
    assert any(p.text == "Sales Overview" and p.heading_level == 2 for p in headings)
    assert any(p.text == "Region Breakdown" and p.heading_level == 2 for p in headings)


def test_parse_extracts_body_text(tmp_path: Path) -> None:
    src = _build_pptx(tmp_path)
    result = PptxParser().parse(src)

    body_texts = [p.text for p in result.paragraphs if p.kind == "text"]
    assert any("Q1 revenue grew 12% YoY" in t for t in body_texts)


def test_parse_extracts_speaker_notes_with_prefix(tmp_path: Path) -> None:
    src = _build_pptx(tmp_path)
    result = PptxParser().parse(src)

    notes_paragraphs = [
        p.text for p in result.paragraphs if p.text.startswith("[Notes]")
    ]
    assert len(notes_paragraphs) == 1
    assert "APAC growth" in notes_paragraphs[0]


def test_parse_extracts_picture_with_sha256(tmp_path: Path) -> None:
    src = _build_pptx(tmp_path)
    result = PptxParser().parse(src)

    assert len(result.embedded_images) == 1
    img = result.embedded_images[0]
    assert img.ext == "png"
    assert len(img.sha256) == 64  # hex digest length
    assert img.image_bytes  # non-empty


def test_parse_extracts_table(tmp_path: Path) -> None:
    src = _build_pptx(tmp_path)
    result = PptxParser().parse(src)

    assert len(result.tables) == 1
    tbl = result.tables[0]
    assert tbl.headers == ["Region", "Revenue"]
    assert tbl.rows == [["APAC", "$1.2M"]]


def test_parse_doc_order_monotonic_across_items(tmp_path: Path) -> None:
    src = _build_pptx(tmp_path)
    result = PptxParser().parse(src)

    all_orders: list[int] = (
        [p.doc_order for p in result.paragraphs]
        + [t.doc_order for t in result.tables]
        + [i.doc_order for i in result.embedded_images]
    )
    assert sorted(all_orders) == list(range(len(all_orders)))


def test_parse_slide_without_title_still_emits_slide_heading(tmp_path: Path) -> None:
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])  # blank layout, no title
    src = tmp_path / "blank.pptx"
    prs.save(src)

    result = PptxParser().parse(src)

    assert not result.parse_failed
    headings = [p for p in result.paragraphs if p.kind == "heading"]
    assert any(p.text == "Slide 1" for p in headings)
    assert not any(p.heading_level == 2 for p in headings)


def test_parse_malformed_file_returns_parse_failed(tmp_path: Path) -> None:
    src = tmp_path / "broken.pptx"
    src.write_bytes(b"not a real pptx")

    result = PptxParser().parse(src)

    assert result.parse_failed
    assert result.parse_error is not None


def test_parser_protocol_doc_format_attribute() -> None:
    assert PptxParser.doc_format == "pptx"
