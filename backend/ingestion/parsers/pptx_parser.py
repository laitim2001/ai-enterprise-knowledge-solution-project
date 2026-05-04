"""python-pptx based .pptx parser (per architecture.md §3.3 + components/C01-ingestion.md §1).

W3 D1 F5 implementation. Mirrors `DoclingDocxParser` design (Parser Protocol
contract; deterministic; never raises) but uses python-pptx because Docling's
.pptx support is less mature than its .docx pipeline and we need explicit
control over slide / placeholder / table / picture extraction.

Slide structure mapping(per slide_based chunk strategy in KbConfig):
- Each slide emits one synthetic level=1 heading "Slide N"(stable anchor for
  F2 chunker section_path)
- Slide title placeholder(if present)= level=2 heading
- Other text frames(body / placeholder / text box)= plain paragraphs
- Tables = `Table` items via doc_order interleave
- Pictures = `EmbeddedImage` via doc_order interleave (SHA256 dedup)
- Speaker notes(if present)appended as `[Notes] ...` plain paragraph

Edge cases:
- Slide with no title placeholder → only "Slide N" heading
- Picture without retrievable blob → skipped (not fatal)
- Malformed / corrupt .pptx → ParserResult(parse_failed=True)

Future polish(W3+ scope):
- Slide layout name → section_path level=2 prefix(currently using title)
- Slide group / section markers from Slides → Sections feature
"""

from __future__ import annotations

import hashlib
import logging
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .base import EmbeddedImage, ParagraphItem, ParserResult, Table

logger = logging.getLogger(__name__)


class PptxParser:
    """python-pptx based .pptx parser. Implements `Parser` Protocol (base.py)."""

    doc_format = "pptx"

    def parse(self, source: Path) -> ParserResult:
        try:
            return self._parse_inner(source)
        except Exception as exc:  # noqa: BLE001 — Parser Protocol MUST NOT raise
            logger.exception("pptx parse failed", extra={"path": str(source)})
            return ParserResult(
                source_path=source,
                doc_format="pptx",
                doc_title=source.stem,
                parse_failed=True,
                parse_error=f"{type(exc).__name__}: {exc}",
            )

    def _parse_inner(self, source: Path) -> ParserResult:
        prs = Presentation(str(source))

        paragraphs: list[ParagraphItem] = []
        embedded_images: list[EmbeddedImage] = []
        tables: list[Table] = []
        doc_order = 0

        for slide_idx, slide in enumerate(prs.slides, start=1):
            paragraphs.append(
                ParagraphItem(
                    text=f"Slide {slide_idx}",
                    kind="heading",
                    doc_order=doc_order,
                    heading_level=1,
                )
            )
            doc_order += 1

            title_text = _extract_title_text(slide)
            title_shape_id = _title_shape_id(slide)
            if title_text:
                paragraphs.append(
                    ParagraphItem(
                        text=title_text,
                        kind="heading",
                        doc_order=doc_order,
                        heading_level=2,
                    )
                )
                doc_order += 1

            for shape in slide.shapes:
                if title_shape_id is not None and shape.shape_id == title_shape_id:
                    continue  # already captured as level=2 heading

                if shape.has_text_frame:
                    text = (shape.text_frame.text or "").strip()
                    if text:
                        paragraphs.append(
                            ParagraphItem(text=text, kind="text", doc_order=doc_order)
                        )
                        doc_order += 1

                elif shape.has_table:
                    raw_rows = [
                        [(cell.text or "").strip() for cell in row.cells]
                        for row in shape.table.rows
                    ]
                    if raw_rows:
                        headers = raw_rows[0]
                        body_rows = raw_rows[1:]
                        tables.append(
                            Table(rows=body_rows, headers=headers, doc_order=doc_order)
                        )
                        doc_order += 1

                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image_record = _try_extract_image(shape, doc_order)
                    if image_record is not None:
                        embedded_images.append(image_record)
                        doc_order += 1

            notes = _extract_notes_text(slide)
            if notes:
                paragraphs.append(
                    ParagraphItem(
                        text=f"[Notes] {notes}", kind="text", doc_order=doc_order
                    )
                )
                doc_order += 1

        return ParserResult(
            source_path=source,
            doc_format="pptx",
            doc_title=source.stem,
            paragraphs=paragraphs,
            embedded_images=embedded_images,
            tables=tables,
        )


def _extract_title_text(slide) -> str:
    try:
        title = slide.shapes.title
    except (AttributeError, ValueError):
        return ""
    if title is None:
        return ""
    try:
        return (title.text or "").strip()
    except AttributeError:
        return ""


def _title_shape_id(slide) -> int | None:
    try:
        title = slide.shapes.title
    except (AttributeError, ValueError):
        return None
    if title is None:
        return None
    return getattr(title, "shape_id", None)


def _try_extract_image(shape, doc_order: int) -> EmbeddedImage | None:
    try:
        blob = shape.image.blob
        ext = shape.image.ext
    except (AttributeError, ValueError, KeyError):
        return None
    if not blob:
        return None
    return EmbeddedImage(
        image_bytes=blob,
        alt_text=getattr(shape, "name", "") or "",
        doc_order=doc_order,
        ext=ext,
        sha256=hashlib.sha256(blob).hexdigest(),
    )


def _extract_notes_text(slide) -> str:
    if not getattr(slide, "has_notes_slide", False):
        return ""
    try:
        return (slide.notes_slide.notes_text_frame.text or "").strip()
    except AttributeError:
        return ""
